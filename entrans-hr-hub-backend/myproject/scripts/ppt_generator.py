import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

def generate_presentation(template_path, excel_path, output_path, user_name, years_of_service):
    df = pd.read_excel(excel_path)
    df.columns = df.columns.str.strip()

    wishes = []
    for _, row in df.iterrows():
        wish = str(row['Wishes']).strip()
        signer = str(row['Name']).strip()
        if not wish or wish.lower() == 'nan' or not signer or signer.lower() == 'nan':
            continue
        wishes.append((wish, signer))

    prs = Presentation(template_path)

    # Place the user name in the correct year slide (index = years_of_service - 1)
    year_slide_idx = int(years_of_service) - 1  # 1-based to 0-based index
    for idx, slide in enumerate(prs.slides):
        if idx == year_slide_idx:
            left = Inches(3)
            top = Inches(3.5)
            width = Inches(4)
            height = Inches(1.2)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            tf.clear()
            p = tf.add_paragraph()
            p.text = user_name
            p.font.size = Pt(24)
            p.font.bold = True
            p.font.name = "Times New Roman"
            p.alignment = PP_ALIGN.CENTER

    def get_blank_layout(prs, year_idx):
      """
      Returns the correct blank layout based on the year index (1-based).
      year_idx: 1, 2, 3, or 4
      """
      if year_idx in [1, 3]:
        return prs.slide_layouts[-2]
      elif year_idx in [2, 4]:
        return prs.slide_layouts[-1]
      else:
        return prs.slide_layouts[-1] 

    # Add one slide per wish (blank layout)
    blank_layout = get_blank_layout(prs, int(years_of_service))
    i = 0
    while i < len(wishes):
        wish_text_1, signer_name_1 = wishes[i]
        wish_text_2, signer_name_2 = None, None

        # Check if next wish exists and both are short
        if (i + 1 < len(wishes)
            and len(wish_text_1) < 120  
            and len(wishes[i + 1][0]) < 120
        ):
            wish_text_2, signer_name_2 = wishes[i + 1]
            i += 2
        else:
            i += 1

        slide = prs.slides.add_slide(blank_layout)

        # First wish
        wish1_top = 1.5
        txBox1 = slide.shapes.add_textbox(Inches(2), Inches(wish1_top), Inches(7), Inches(1))
        tf1 = txBox1.text_frame
        tf1.clear()
        tf1.word_wrap = True
        p1 = tf1.add_paragraph()
        p1.font.name = "Times New Roman"
        p1.text = wish_text_1
        p1.font.size = Pt(18)
        p1.font.bold = False
        p1.font.italic = False
        p1.alignment = 1  # Center alignment

        # Dynamic Y for signer based on wish length
        signer1_offset = 0.5 if len(wish_text_1) < 80 else 1.0 + ((len(wish_text_1) - 80) // 80) * 0.3
        signer1_top = wish1_top + signer1_offset
        signer_box1 = slide.shapes.add_textbox(Inches(3), Inches(signer1_top), Inches(5), Inches(0.6))
        signer_tf1 = signer_box1.text_frame
        signer_tf1.clear()
        signer_tf1.word_wrap = True
        signer_p1 = signer_tf1.add_paragraph()
        signer_p1.font.name = "Times New Roman"
        signer_p1.font.color.rgb = RGBColor(255, 0, 0)
        signer_p1.text = f"- {signer_name_1}"
        signer_p1.font.size = Pt(18)
        signer_p1.font.bold = True
        signer_p1.font.italic = True
        signer_p1.alignment = 3  # Right alignment

        # Second wish (if present)
        if wish_text_2:
            wish2_top = signer1_top + 0.8  # Add space after first signer
            txBox2 = slide.shapes.add_textbox(Inches(2), Inches(wish2_top), Inches(7), Inches(1))
            tf2 = txBox2.text_frame
            tf2.clear()
            tf2.word_wrap = True
            p2 = tf2.add_paragraph()
            p2.font.name = "Times New Roman"
            p2.text = wish_text_2
            p2.font.size = Pt(18)
            p2.font.bold = False
            p2.font.italic = False
            p2.alignment = 1  # Center alignment

            # Dynamic Y for second signer
            signer2_offset = 0.5 if len(wish_text_2) < 80 else 1.0 + ((len(wish_text_2) - 80) // 80) * 0.3
            signer2_top = wish2_top + signer2_offset
            signer_box2 = slide.shapes.add_textbox(Inches(3), Inches(signer2_top), Inches(5), Inches(0.6))
            signer_tf2 = signer_box2.text_frame
            signer_tf2.clear()
            signer_tf2.word_wrap = True
            signer_p2 = signer_tf2.add_paragraph()
            signer_p2.font.name = "Times New Roman"
            signer_p2.font.color.rgb = RGBColor(255, 0, 0)
            signer_p2.text = f"- {signer_name_2}"
            signer_p2.font.size = Pt(18)
            signer_p2.font.bold = True
            signer_p2.font.italic = True
            signer_p2.alignment = 3  # Right alignment

    # Add a Thank You slide at the end
    thank_slide = prs.slides.add_slide(blank_layout)
    txBox = thank_slide.shapes.add_textbox(Inches(3.2), Inches(2), Inches(6), Inches(2))
    tf = txBox.text_frame
    tf.clear()
    p = tf.add_paragraph()
    p.text = "Thank you"
    p.font.size = Pt(64)
    p.font.name = "Times New Roman"
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 0, 0)  # Red color
    p.font.italic = True
    p.alignment = 1  # Center

    # --- Remove all template slides except the relevant year slide (first 6 slides) ---
    slides_to_remove = [i for i in range(6) if i != year_slide_idx]
    for i in sorted(slides_to_remove, reverse=True):
        xml_slides = prs.slides._sldIdLst
        xml_slides.remove(xml_slides[i])

    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

